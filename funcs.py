import json
from random import randint

from docxtpl import DocxTemplate

import xlsxwriter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor



def read_data():
    with open("data.json", "r", encoding="utf-8") as file:
        data = json.load(file)
        return data


def create_json():
    with open("data.json", "w") as file:
        data = {"cinemas": {}, "films": {}}
        json.dump(data, file)

def init_data(ts):
    ts.add_cinema("Кинотеатр-1")
    ts.add_film("Фильм-1", 123)
    ts.add_film("Фильм-2", 114)
    ts.add_hall("Кинотеатр-1", "Первый")
    ts.add_session("Кинотеатр-1", "Первый", "Фильм-1", "07:00")
    ts.add_session("Кинотеатр-1", "Первый", "Фильм-2", "10:00")
    ts.add_session("Кинотеатр-1", "Первый", "Фильм-2", "12:00")

def print_dict(dct):
    for key, value in dct.items():
        print(key, value)


def create_report_timetable(month_name, data):
    doc = DocxTemplate("templates/template_report_timetable.docx")
    context = {"month_name": month_name,
               "data": data
               }
    doc.render(context)
    file_name = f"reports/Расписание сеансов за {month_name}.docx"
    doc.save(file_name)

def create_report_grafik(data):
    workbook = xlsxwriter.Workbook('reports/график загруженности кинотеатров.xlsx')
    worksheet = workbook.add_worksheet()

    worksheet.write_column('A1', data)
    chart = workbook.add_chart({'type': 'column'})
    chart.set_title({'name': 'График загруженности кинотеатров'})
    chart.set_x_axis({'name': 'Часы'})
    chart.set_y_axis({'name': 'Доля проданых билетов'})
    chart.set_size({'width': 720, 'height': 576})
    chart.set_style(37)

    n = len(data)//25
    for i in range(n):
        chart.add_series({'values': f'=Sheet1!A{i * 25 + 2}:A{i * 25 + 25}',
                          'name': f'=Sheet1!A{i * 25 + 1}'})

    worksheet.insert_chart('C1', chart)
    workbook.close()

def create_report_presentation(films, data):

    epitets = ["самый кассовый фильм года", "остросюжетный фильм",
              "во всех кинотеатрах нашего старого знакомого фильм", "самый убойный фильм этого лета",
              "лучший фильм от знаменитого режиссера"]
    prs = Presentation()

    for film in films:

        title_slide_layout = prs.slide_layouts[8]

        slide = prs.slides.add_slide(title_slide_layout)

        title = slide.placeholders[0]
        num = randint(0, len(epitets) - 1)
        title.text = f'Смотрите {epitets[num]} {film}'

        pic = slide.placeholders[1]
        film_name = "_".join(film.split())
        try:
            pic.insert_picture(f'images/{film_name}.jpg')
        except:
            pass

        lst_cinemas = data[film]
        text = slide.placeholders[2]
        text.text = f"Премьеры в кинотеатрах: {', '.join(lst_cinemas)}"



    prs.save('reports/Буклеты фильмов.pptx')